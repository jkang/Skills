from pptx import Presentation
from pptx.util import Inches, Pt

def create_test_ppt(path):
    prs = Presentation()
    
    # Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Existing Content Test"
    subtitle.text = "This should be restyled but preserved."
    
    # Content Slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Slide Header"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Bullet point 1"
    p = tf.add_paragraph()
    p.text = "Bullet point 2 with some detail"
    
    # Table Slide
    slide_layout = prs.slide_layouts[5] # Blank
    slide = prs.slides.add_slide(slide_layout)
    rows = 2
    cols = 2
    left = top = Inches(2)
    width = Inches(4)
    height = Inches(1)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    table.cell(0, 0).text = "Header A"
    table.cell(0, 1).text = "Header B"
    table.cell(1, 0).text = "Data 1"
    table.cell(1, 1).text = "Data 2"

    prs.save(path)
    print(f"Test PPT created at {path}")

if __name__ == "__main__":
    create_test_ppt("test_target.pptx")
