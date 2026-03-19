#!/usr/bin/env python3
"""
InspireTemplate Style Transfer - v3.1
=======================================
核心策略：
  1. 结构性页面 (封面/议程/章节分隔页/结尾页):
     → 直接从 InspireTemplate.pptx 复制对应幻灯片 XML
       再将源 PPT 的关键文本注入到复制的幻灯片中
  2. 内容页 / 互动页:
     → 保留源幻灯片的所有形状和布局
       应用 Inspire 样式：
         - 背景：白色 (#FFFFFF)
         - 标题：思源黑体+粗体+#1B2B47
         - 正文：思源黑体+#1B2B47
         - 表格：表头深蓝背景+白字

InspireTemplate.pptx 幻灯片映射 (已知，固定):
  index 0: 说明页（跳过）
  index 1: Cover 封面
  index 4: Agenda/Contents 目录
  index 5: Section Divider 章节封面（无图）
  index 6: Section Divider with Image（备用）
  index 10: Ending 结束页

Usage:
  python style_transfer.py <target.pptx>
  python style_transfer.py <target.pptx> --template InspireTemplate.pptx
"""

import os, sys, copy, json, argparse, traceback
from lxml import etree
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn


# ─── Brand Colors (Inspire Design System) ────────────────────────────────────
STARRY_BLUE    = "#1B2B47"   # 星空蓝 - 主品牌色：背景、标题
CREATIVE_BLUE  = "#4A9FD8"   # 创想蓝 - 强调色：数字序号、bullet点
TECH_GRAY      = "#F5F7FA"   # 科技灰 - 内容页背景、卡片背景
WHITE          = "#FFFFFF"   # 白色    - 内容页主背景/深色背景文字
TEXT_PRIMARY   = "#1B2B47"   # 正文主色（深蓝灰，非纯黑）
TEXT_SECONDARY = "#64748B"   # 次要文字（副标题、图注）
BORDER_COLOR   = "#E2E8F0"   # 边框/分割线

# InspireTemplate.pptx 中固定的幻灯片索引映射
TEMPLATE_SLIDE_MAP = {
    "cover":       1,   # Inspire PPT Template Cover
    "agenda":      4,   # CONTENTS 目录页
    "section":     5,   # Slide Separator 章节封面（无图）
    "section_img": 6,   # Slide Separator With Image
    "ending":      10,  # We look forward to working with you
}

# 中文字体 (源文件往往有思源黑体/苹方，这里用微软雅黑作通用回退)
HEADING_FONT = "Microsoft YaHei"
BODY_FONT    = "Microsoft YaHei"


# ─── Color helpers ────────────────────────────────────────────────────────────
def hex_to_rgb(h: str) -> RGBColor:
    h = h.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def set_slide_solid_bg(slide, hex_color: str):
    """Give a slide a solid color background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(hex_color)


# ─── Text helpers ─────────────────────────────────────────────────────────────
def get_slide_full_text(slide) -> str:
    parts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                parts.append(para.text)
    return ' '.join(parts)


def extract_texts_by_importance(slide) -> list:
    """
    Return all non-empty paragraph texts from slide,
    ordered: largest font size → smallest.
    """
    items = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            max_sz = 0
            is_bold = False
            for run in para.runs:
                if run.font.size and run.font.size.pt > max_sz:
                    max_sz = run.font.size.pt
                if run.font.bold:
                    is_bold = True
            items.append((max_sz, is_bold, text))
    items.sort(key=lambda x: (x[0], x[1]), reverse=True)
    return [item[2] for item in items]


# ─── Slide classification (for SOURCE PPT slides) ────────────────────────────
def classify_slide(slide, index: int, total: int) -> str:
    """
    Classify a source slide into one of:
    'cover', 'agenda', 'section', 'ending', 'interactive', 'content'
    
    Based on analysis of the 55-page reference PDF (AI商业落地实战):
    - Page 1      = cover
    - Page 2      = agenda (AGENDA / 议程)
    - PART XX     = section (dark gradient bg, chapter divider)
    - 🎯          = interactive
    - Last page   = ending
    - All others  = content
    """
    if index == 0:
        return 'cover'

    text = get_slide_full_text(slide)
    text_stripped = text.strip()
    text_lower = text_stripped.lower()
    lines = [l.strip() for l in text_stripped.split() if l.strip()]
    first_words = text_stripped[:60].lower()

    # Last slide: ending page
    if index == total - 1 and len(text_stripped) < 400:
        return 'ending'

    # Agenda (appears near the start, contains agenda/目录/议程)
    agenda_kw = ['agenda', '议程', '目录', 'contents', '大纲', 'outline']
    if index <= 3 and any(kw in text_lower for kw in agenda_kw):
        return 'agenda'

    # Section dividers: "PART XX" pattern (very characteristic in this template)
    # These are short slides with dark gradient background
    import re
    if re.search(r'\bpart\s+\d+', text_lower):
        # Only if the text is relatively short (it's a divider, not a content page mentioning PART)
        if len(text_stripped) < 400:
            return 'section'

    # Also catch PART SUMMARY pages (they're styled as content, not section)
    # so don't classify them as section
    
    # Interactive pages: start with 🎯 or contain strong interactive keywords
    if '🎯' in text or '互动环节' in text or '实操引导' in text:
        return 'interactive'
    
    interactive_kw = ['小组讨论', '请思考', '课堂活动', '让我们一起']
    if any(kw in text for kw in interactive_kw):
        return 'interactive'

    return 'content'


# ─── Slide cloning from template ─────────────────────────────────────────────
def clone_template_slide(template_prs: Presentation, slide_index: int,
                          target_prs: Presentation):
    """
    Deep-clone a slide from template_prs[slide_index] into target_prs.
    Returns the new slide.
    """
    try:
        template_slide = template_prs.slides[slide_index]
    except IndexError:
        print(f"  [warn] Template slide index {slide_index} out of range, skipping clone")
        return None

    # Add blank slide to target
    blank_layout = target_prs.slide_layouts[6]
    new_slide = target_prs.slides.add_slide(blank_layout)

    # Replace spTree (all shapes)
    tpl_spTree = template_slide.shapes._spTree
    new_spTree = new_slide.shapes._spTree
    for child in list(new_spTree):
        new_spTree.remove(child)
    for child in tpl_spTree:
        new_spTree.append(copy.deepcopy(child))

    # Copy background
    tpl_bg = template_slide.background._element
    new_bg = new_slide.background._element
    for child in list(new_bg):
        new_bg.remove(child)
    for child in tpl_bg:
        new_bg.append(copy.deepcopy(child))

    return new_slide


def get_text_shapes_sorted_by_size(slide) -> list:
    """
    Return list of (estimated_size, shape) tuples sorted largest-first.
    Only includes shapes that have non-empty text.
    """
    result = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if not shape.text_frame.text.strip():
            continue
        max_sz = 0
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size and run.font.size.pt > max_sz:
                    max_sz = run.font.size.pt
        result.append((max_sz, shape))
    result.sort(key=lambda x: x[0], reverse=True)
    return result


def inject_texts_into_cloned_slide(cloned_slide, texts: list):
    """
    Replace text in the TOP-N largest text shapes of a cloned template slide
    with the provided texts (ordered by importance).
    Each shape's first run gets the corresponding text.
    """
    text_shapes = get_text_shapes_sorted_by_size(cloned_slide)

    for i, (_, shape) in enumerate(text_shapes):
        if i >= len(texts) or not texts[i]:
            break
        new_text = texts[i]
        tf = shape.text_frame
        # Find first paragraph with runs
        set_done = False
        for para in tf.paragraphs:
            if para.runs and not set_done:
                para.runs[0].text = new_text
                for run in para.runs[1:]:
                    run.text = ''
                set_done = True
            elif set_done:
                # Clear remaining paragraphs
                for run in para.runs:
                    run.text = ''


# ─── Content slide styling ────────────────────────────────────────────────────
def is_title_shape(shape) -> bool:
    """Heuristic: is this shape likely a slide title?"""
    if not shape.has_text_frame:
        return False
    # Positional: near the top of the slide
    if shape.top is not None and shape.top < Emu(1_800_000):
        return True
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size and run.font.size.pt >= 24:
                return True
            if run.font.bold and len(run.text.strip()) < 70:
                return True
    return False


def style_run(run, is_title: bool = False):
    """Apply Inspire font + color to a single text run."""
    try:
        run.font.name = HEADING_FONT if is_title else BODY_FONT
        if is_title:
            run.font.bold = True
            run.font.color.rgb = hex_to_rgb(STARRY_BLUE)
        else:
            run.font.color.rgb = hex_to_rgb(TEXT_PRIMARY)
    except Exception:
        pass


def style_table(table):
    """Apply Inspire table styling: header row = dark bg + white bold text."""
    for row_idx, row in enumerate(table.rows):
        is_header = (row_idx == 0)
        for col_idx, cell in enumerate(row.cells):
            try:
                if is_header:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = hex_to_rgb(STARRY_BLUE)
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.color.rgb = hex_to_rgb(WHITE)
                            run.font.name = HEADING_FONT
                            run.font.bold = True
                else:
                    # First column: slightly accented
                    text_color = STARRY_BLUE if col_idx == 0 else TEXT_PRIMARY
                    bold = col_idx == 0
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.color.rgb = hex_to_rgb(text_color)
                            run.font.name = BODY_FONT
                            run.font.bold = bold
            except Exception:
                pass


def style_shape(shape):
    """Apply Inspire styles to a single shape (text + table)."""
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                style_shape(child)
            return

        if shape.has_text_frame:
            title = is_title_shape(shape)
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    style_run(run, is_title=title)

        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            style_table(shape.table)

    except Exception as e:
        print(f"  [warn] Cannot style shape '{shape.name}': {e}")


def style_content_slide(slide, is_interactive: bool = False):
    """Apply full Inspire styling to a content/interactive slide."""
    bg_color = "#FFFDF5" if is_interactive else WHITE
    set_slide_solid_bg(slide, bg_color)
    for shape in slide.shapes:
        style_shape(shape)


# ─── Copy source shapes into a new slides ────────────────────────────────────
def copy_shapes_from_source(source_slide, target_slide):
    """Deep-copy all shapes from source_slide to target_slide."""
    src_spTree = source_slide.shapes._spTree
    tgt_spTree = target_slide.shapes._spTree
    for child in list(tgt_spTree):
        tgt_spTree.remove(child)
    for child in src_spTree:
        tgt_spTree.append(copy.deepcopy(child))


# ─── Main Transformer ─────────────────────────────────────────────────────────
class StyleTransfer:
    def __init__(self, template_path: str, style_config_path: str = None):
        self.template_path = template_path
        self.template_prs  = Presentation(template_path)
        n = len(self.template_prs.slides)
        print(f"[template] {n} slides loaded from {os.path.basename(template_path)}")
        for k, idx in TEMPLATE_SLIDE_MAP.items():
            if idx < n:
                tpl_text = get_slide_full_text(self.template_prs.slides[idx])[:60]
                print(f"  → [{k}] slide {idx+1}: \"{tpl_text}\"")
            else:
                print(f"  → [{k}] slide {idx+1}: OUT OF RANGE (template only has {n} slides)")

    def _build_structural_slide(self, source_slide, slide_type: str,
                                 output_prs: Presentation):
        """Clone template slide for structural types and inject source texts."""
        idx = TEMPLATE_SLIDE_MAP.get(slide_type)
        if idx is None:
            return None
        n_tpl = len(self.template_prs.slides)
        if idx >= n_tpl:
            print(f"  [warn] Template slide {idx} for '{slide_type}' not available")
            return None

        # --- Clone the template slide ---
        cloned = clone_template_slide(self.template_prs, idx, output_prs)
        if cloned is None:
            return None

        # --- Extract key texts from source ---
        key_texts = extract_texts_by_importance(source_slide)
        preview = key_texts[:2] if key_texts else []
        print(f"  [clone] '{slide_type}' → injecting {len(key_texts)} text(s): {preview}")

        inject_texts_into_cloned_slide(cloned, key_texts)
        return cloned

    def transform(self, target_path: str) -> str:
        print(f"\n{'='*60}")
        print(f"[InspireTransfer v3.1]  {os.path.basename(target_path)}")
        print(f"{'='*60}")

        source_prs = Presentation(target_path)
        output_prs = Presentation()

        # Match slide dimensions
        output_prs.slide_width  = source_prs.slide_width
        output_prs.slide_height = source_prs.slide_height

        total = len(source_prs.slides)
        print(f"[source]  {total} slides\n")

        for idx, src_slide in enumerate(source_prs.slides):
            slide_type = classify_slide(src_slide, idx, total)
            preview = get_slide_full_text(src_slide)[:70].replace('\n', ' ')
            print(f"  [{idx+1:02d}/{total}] → {slide_type:12s}  \"{preview}\"")

            if slide_type in ('cover', 'agenda', 'section', 'ending'):
                # Clone the matching InspireTemplate slide
                new_slide = self._build_structural_slide(src_slide, slide_type, output_prs)
                if new_slide is None:
                    # Fallback: treat as content
                    print(f"  [fallback] treating as content")
                    blank = output_prs.slide_layouts[6]
                    new_slide = output_prs.slides.add_slide(blank)
                    copy_shapes_from_source(src_slide, new_slide)
                    style_content_slide(new_slide)
            else:
                # Content or interactive: copy shapes + apply Inspire styles
                blank = output_prs.slide_layouts[6]
                new_slide = output_prs.slides.add_slide(blank)
                copy_shapes_from_source(src_slide, new_slide)
                style_content_slide(new_slide, is_interactive=(slide_type == 'interactive'))

        output_path = os.path.splitext(target_path)[0] + "_inspired.pptx"
        output_prs.save(output_path)
        print(f"\n✅  Saved → {output_path}\n")
        return output_path


# ─── CLI entry point ──────────────────────────────────────────────────────────
def main():
    parser = argparse.ArgumentParser(
        description="InspireTemplate Style Transfer v3.1"
    )
    parser.add_argument("target_ppt",
        help="Source PPT file to transform")
    parser.add_argument("--template",
        default=os.path.join(os.path.dirname(__file__), "..", "InspireTemplate.pptx"),
        help="Path to InspireTemplate.pptx")
    parser.add_argument("--style",
        default=os.path.join(os.path.dirname(__file__), "..", "pptstyle.json"),
        help="Path to pptstyle.json (informational)")
    args = parser.parse_args()

    if not os.path.exists(args.target_ppt):
        print(f"❌ Error: '{args.target_ppt}' not found")
        sys.exit(1)
    if not os.path.exists(args.template):
        print(f"❌ Error: template '{args.template}' not found")
        sys.exit(1)

    try:
        transfer = StyleTransfer(args.template, args.style)
        transfer.transform(args.target_ppt)
    except Exception:
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()
