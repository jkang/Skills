from pptx import Presentation
import sys, copy
sys.path.append('scripts')
from style_transfer import *

host_prs = Presentation('InspireTemplate.pptx')
new_slide = clone_template_slide(host_prs, 1) # Cover

inject_texts(new_slide, ['AI2.0'])

print("\nAfter injection (memory XML):")
for p in new_slide._element.xpath('.//a:p/a:r/a:t'):
    print('  text=', p.text)

host_prs.save('test_xml.pptx')

prs2 = Presentation('test_xml.pptx')
s2 = prs2.slides[-1]
print("\nAfter save (loaded XML):")
for p in s2._element.xpath('.//a:p/a:r/a:t'):
    print('  text=', p.text)
