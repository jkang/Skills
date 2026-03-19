from pptx import Presentation
import copy
tpl_prs = Presentation('InspireTemplate.pptx')
new_slide = tpl_prs.slides.add_slide(tpl_prs.slide_layouts[6])
tgt_tree = new_slide._element.spTree
for child in list(tgt_tree): tgt_tree.remove(child)
for child in tpl_prs.slides[1]._element.spTree:
    tgt_tree.append(copy.deepcopy(child))

new_val = 'NEW TITLE TEST'
for shape in new_slide.shapes:
    if shape.has_text_frame:
        set_done = False
        for p in shape.text_frame.paragraphs:
            if p.runs and not set_done:
                p.runs[0].text = new_val
                for r in p.runs[1:]:
                    r.text = ''
                set_done = True
            elif set_done:
                for r in p.runs:
                    r.text = ''
tpl_prs.save('test_inject.pptx')

prs2 = Presentation('test_inject.pptx')
s2 = prs2.slides[-1]
print('Saved shape text:', [s.text_frame.text for s in s2.shapes if s.has_text_frame])
