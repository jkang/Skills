from pptx import Presentation
import os

def debug_pptx(path):
    print(f"--- Debugging {os.path.basename(path)} ---")
    prs = Presentation(path)
    
    # Check Theme fonts
    master = prs.slide_masters[0]
    print(f"Master Name: {master.name}")
    
    # Check placeholders on first slide
    if len(prs.slides) > 0:
        slide = prs.slides[0]
        print(f"Slide 1 shapes: {len(slide.shapes)}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        font_name = run.font.name
                        font_size = run.font.size.pt if run.font.size else "None"
                        print(f"Shape: {shape.name}, Text: {run.text[:20]}..., Font: {font_name}, Size: {font_size}")

if __name__ == "__main__":
    template = "/Users/superkkk/MyCoding/Skills/myskills/WorkSkills/InspirePPTSkills/InspireTemplate.pptx"
    target = "/Users/superkkk/MyCoding/Skills/myskills/WorkSkills/InspirePPTSkills/AI 2.0 产品流程和实践.pptx"
    debug_pptx(template)
    print("\n")
    debug_pptx(target)
