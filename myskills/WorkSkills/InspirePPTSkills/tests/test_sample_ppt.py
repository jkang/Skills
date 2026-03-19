import traceback
from pptx import Presentation

try:
    prs = Presentation('../SamplePPT.pptx')
    print(f"SamplePPT.pptx loaded with {len(prs.slides)} slides.")
    
    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i+1} ---")
        txt = ""
        try:
            for s in slide.shapes:
                try:
                    if s.has_text_frame:
                        txt += s.text_frame.text.replace('\n', ' ')[:50] + " | "
                except Exception as e:
                    pass
            print(f"  Texts: {txt}")
        except Exception as e:
            print(f"  Error reading shapes: {e}")
            
except Exception as e:
    print(f"Fatal error: {e}")
