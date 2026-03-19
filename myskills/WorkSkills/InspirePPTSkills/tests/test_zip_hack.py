from pptx import Presentation
import sys, copy, zipfile, os, re
sys.path.append('scripts')
from style_transfer import *

src_prs = Presentation('AI 2.0 产品流程和实践.pptx')
host_prs = Presentation('InspireTemplate.pptx')

slide = src_prs.slides[0]
texts = extract_texts_by_size(slide)

# Clone and inject
new_slide = clone_template_slide(host_prs, 1) # Cover
inject_texts(new_slide, texts)

# DO NOT delete slides. Just save.
host_prs.save('test_zip_hack_temp.pptx')

# Unzip and modify presentation.xml
with zipfile.ZipFile('test_zip_hack_temp.pptx', 'r') as zin:
    with zipfile.ZipFile('test_zip_hack.pptx', 'w') as zout:
        for item in zin.infolist():
            content = zin.read(item.filename)
            if item.filename == 'ppt/presentation.xml':
                xml = content.decode('utf-8')
                matches = list(re.finditer(r'<p:sldId id="\d+" r:id="rId\d+"/>', xml))
                if len(matches) >= 11:
                    to_remove = matches[:11]
                    for m in to_remove:
                        xml = xml.replace(m.group(0), '')
                content = xml.encode('utf-8')
            zout.writestr(item, content)

prs2 = Presentation('test_zip_hack.pptx')
s2 = prs2.slides[0]
print('Saved text in zip-hacked PPT: ', [s.text_frame.text.strip() for s in s2.shapes if s.has_text_frame])
