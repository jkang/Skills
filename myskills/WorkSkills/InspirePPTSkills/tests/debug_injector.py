from pptx import Presentation
import sys, os
sys.path.append('scripts')
from style_transfer import *

src_prs = Presentation('AI 2.0 产品流程和实践.pptx')
host_prs = Presentation('InspireTemplate.pptx')

slide = src_prs.slides[0]
texts = extract_texts_by_size(slide)
print('Source texts:', texts)

new_slide = clone_template_slide(host_prs, 1) # Cover

print("Before injection:")
for s in new_slide.shapes:
    if s.has_text_frame:
        print("  Found text-frame:", [p.text for p in s.text_frame.paragraphs])

inject_texts(new_slide, texts)

print("After injection:")
for s in new_slide.shapes:
    if s.has_text_frame:
        print("  Found text-frame:", [p.text for p in s.text_frame.paragraphs])
