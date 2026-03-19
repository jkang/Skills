from pptx import Presentation
import sys, copy, re
sys.path.append('scripts')
from style_transfer import *

src_prs = Presentation('AI 2.0 产品流程和实践.pptx')
host_prs = Presentation('InspireTemplate.pptx')

slide = src_prs.slides[0]
texts = extract_texts_by_size(slide)

# Clone
new_slide = clone_template_slide(host_prs, 1) # Cover

# DIRECT XML REPLACEMENT
xml_str = new_slide._element.xml
print('Before XML replace, contains Inspire:', 'Inspire PPT Template' in xml_str)

# We know the template has "Inspire PPT Template" and "Cover"
xml_str = xml_str.replace('Inspire PPT Template ', texts[0])
xml_str = xml_str.replace('Cover', texts[1] + ' ' + texts[2])

from pptx.oxml import parse_xml
new_slide._element.getparent().replace(new_slide._element, parse_xml(xml_str))

# Delete front slides
tpl_count = 11
sldIdLst = host_prs.slides._sldIdLst
for sldId in list(sldIdLst)[:tpl_count]:
    sldIdLst.remove(sldId)

host_prs.save('test_xml_replace.pptx')

prs2 = Presentation('test_xml_replace.pptx')
s2 = prs2.slides[0]
print('Saved text in XML replacing PPT: ', [s.text_frame.text.strip() for s in s2.shapes if s.has_text_frame])
