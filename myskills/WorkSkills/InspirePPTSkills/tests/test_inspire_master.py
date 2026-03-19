from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

prs = Presentation('InspireTemplate.pptx')
master = prs.slide_master

print("=== INSPIRE MASTER SHAPES ===")
for shape in master.shapes:
    print(f"Shape: {shape.name}, Type: {shape.shape_type}")
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        print(f"  [PICTURE] {shape.width.pt} x {shape.height.pt} at {shape.left.pt}, {shape.top.pt}")
    if shape.has_text_frame:
        print(f"  [TEXT] {shape.text_frame.text}")

print("\n=== INSPIRE LAYOUT 3 SHAPES (Content Slide) ===")
layout = prs.slide_layouts[2]
for shape in layout.shapes:
    print(f"Shape: {shape.name}, Type: {shape.shape_type}")
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        print(f"  [PICTURE] {shape.width.pt} x {shape.height.pt} at {shape.left.pt}, {shape.top.pt}")
    if shape.has_text_frame:
        print(f"  [TEXT] {shape.text_frame.text}")
