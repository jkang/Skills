from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

prs = Presentation('AI 2.0 产品流程和实践_inspired.pptx')
slide = prs.slides[6]  # Slide 7 (0-indexed)

def inspect_shape(shape, indent=""):
    print(f"{indent}Shape: {shape.name}, Type: {shape.shape_type}")
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        print(f"{indent}  [GROUP] containing {len(shape.shapes)} shapes:")
        for child in shape.shapes:
            inspect_shape(child, indent + "    ")
    
    if hasattr(shape, "fill") and shape.fill.type == 1:
        if hasattr(shape.fill.fore_color, 'rgb') and shape.fill.fore_color.rgb:
            print(f"{indent}  Fill RGB: {shape.fill.fore_color.rgb}")
        elif hasattr(shape.fill.fore_color, 'theme_color'):
            print(f"{indent}  Fill Theme Color: {shape.fill.fore_color.theme_color}")

print("=== SLIDE 7 SHAPES ===")
for shape in slide.shapes:
    inspect_shape(shape)

print("\n=== SLIDE 7 FOOTER CHECK ===")
for shape in slide.shapes:
    if shape.has_text_frame:
        print(f"Text: {shape.text_frame.text.replace(chr(11), ' ')}")

print("\n=== ALL SLIDES LOGO CHECK ===")
for idx, s in enumerate(prs.slides):
    for shape in s.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            print(f"Slide {idx+1} Picture: {shape.name} ({shape.width.pt} x {shape.height.pt} at {shape.left.pt}, {shape.top.pt})")

print("\n=== LAYOUT LOGO CHECK ===")
for i, layout in enumerate(prs.slide_layouts):
    for shape in layout.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            print(f"Layout {i+1} Picture: {shape.name} ({shape.width.pt} x {shape.height.pt})")
