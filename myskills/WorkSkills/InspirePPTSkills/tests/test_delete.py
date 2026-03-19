from pptx import Presentation
import sys, copy
sys.path.append('scripts')
from style_transfer import *

src_prs = Presentation('AI 2.0 产品流程和实践.pptx')
host_prs = Presentation('InspireTemplate.pptx')

slide = src_prs.slides[0]
texts = extract_texts_by_size(slide)

# Clone and inject
new_slide = clone_template_slide(host_prs, 1) # Cover
inject_texts(new_slide, texts)

# Delete front slides
tpl_count = 11
sldIdLst = host_prs.slides._sldIdLst
for sldId in list(sldIdLst)[:tpl_count]:
    sldIdLst.remove(sldId)

for s in host_prs.slides:
    print('Before save, slide text:', [sh.text_frame.text.strip() for sh in s.shapes if sh.has_text_frame])

host_prs.save('test_delete.pptx')

prs2 = Presentation('test_delete.pptx')
s2 = prs2.slides[0]
print('Saved text in deleted+saved PPT: ', [s.text_frame.text.strip() for s in s2.shapes if s.has_text_frame])
