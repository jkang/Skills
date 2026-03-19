from pptx import Presentation
import os

def check_inspired_ppt(path):
    print(f"--- Checking {os.path.basename(path)} ---")
    prs = Presentation(path)
    
    if len(prs.slides) > 0:
        slide = prs.slides[0]
        print(f"Slide 1 shapes: {len(slide.shapes)}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        font_name = run.font.name
                        print(f"Shape: {shape.name}, Text: {run.text[:20]}..., Font: {font_name}")

if __name__ == "__main__":
    target = "/Users/superkkk/MyCoding/Skills/myskills/WorkSkills/InspirePPTSkills/AI 2.0 产品流程和实践_inspired.pptx"
    check_inspired_ppt(target)
