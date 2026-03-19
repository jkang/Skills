import io
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import copy

prs = Presentation('InspireTemplate.pptx')
out_prs = Presentation('AI 2.0 替换主题后.pptx')  # Use the theme-replaced one!

tpl_slide = prs.slides[1] # Cover
tgt_slide = out_prs.slides[0]

# Clear tgt shapes
for s in list(tgt_slide.shapes):
    s._element.getparent().remove(s._element)

# 1. Copy Layout shapes (Skipping Placeholders)
layout = tpl_slide.slide_layout
for shape in layout.shapes:
    if shape.is_placeholder:
        continue
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        blob = shape.image.blob
        tgt_slide.shapes.add_picture(io.BytesIO(blob), shape.left, shape.top, shape.width, shape.height)
    else:
        # deepcopy vector shapes
        tgt_slide._element.spTree.append(copy.deepcopy(shape._element))

# 2. Copy Slide shapes
for shape in tpl_slide.shapes:
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        blob = shape.image.blob
        tgt_slide.shapes.add_picture(io.BytesIO(blob), shape.left, shape.top, shape.width, shape.height)
    else:
        tgt_slide._element.spTree.append(copy.deepcopy(shape._element))

out_prs.save('flatten_test.pptx')
print('Saved flatten_test.pptx')
