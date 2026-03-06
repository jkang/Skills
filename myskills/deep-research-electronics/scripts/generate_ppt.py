import json
import sys
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_ppt(data, output_path):
    prs = Presentation()
    
    # Use a blank layout for maximum control
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # 5 companies, 5 columns
    companies = data.get("companies", [])
    num_cols = len(companies)
    if num_cols == 0:
        print("No company data provided.")
        return

    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    margin = Inches(0.5)
    col_width = (slide_width - (2 * margin)) / num_cols
    
    # Collect all sources for the notes section
    all_sources = []
    
    for i, company in enumerate(companies):
        left = margin + (i * col_width)
        top = margin
        
        # Company Name (Header)
        txBox = slide.shapes.add_textbox(left, top, col_width, Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = company.get("name", "Unknown")
        p.font.bold = True
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.CENTER
        
        # Movements
        top += Inches(0.6)
        txBox_content = slide.shapes.add_textbox(left, top, col_width, slide_height - margin - top)
        tf_content = txBox_content.text_frame
        tf_content.word_wrap = True
        
        for movement in company.get("movements", []):
            p = tf_content.add_paragraph()
            p.text = f"• {movement}"
            p.font.size = Pt(10)
            p.space_after = Pt(5)
            
        # Collect sources
        sources = company.get("sources", [])
        if sources:
            all_sources.append(f"{company['name']} Sources:")
            all_sources.extend(sources)
            all_sources.append("") # spacer

    # Add Notes
    if all_sources:
        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        notes_text_frame.text = "\n".join(all_sources)

    prs.save(output_path)
    print(f"PPT saved to {output_path}")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate a minimalist multi-column PPT report.")
    parser.add_argument("json_file", help="Path to the JSON file containing research data.")
    parser.add_argument("output_ppt", help="Path to save the generated PPT file.")
    
    args = parser.parse_args()
    
    try:
        with open(args.json_file, 'r', encoding='utf-8') as f:
            data = json.load(f)
        create_ppt(data, args.output_ppt)
    except Exception as e:
        print(f"Error: {e}")
        sys.exit(1)
